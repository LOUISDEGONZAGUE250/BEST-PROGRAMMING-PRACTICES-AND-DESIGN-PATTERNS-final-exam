import base64
import urllib.request
import urllib.error
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def get_mermaid_image(mermaid_text, output_path):
    try:
        # Mermaid.ink expects base64 encoded string
        # Actually it expects a JSON object {"code": "...", "mermaid": {"theme": "default"}}
        # But for simple img API: base64 of the text works if formatted correctly.
        # Let's use the standard base64 encoding
        encoded = base64.b64encode(mermaid_text.encode('utf-8')).decode('ascii')
        url = f"https://mermaid.ink/img/{encoded}"
        
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req) as response:
            with open(output_path, 'wb') as out_file:
                out_file.write(response.read())
        return True
    except Exception as e:
        print(f"Error fetching image for {output_path}: {e}")
        return False

# Initialize presentation
prs = Presentation()

# Layouts
TITLE_SLIDE_LAYOUT = 0
BULLET_SLIDE_LAYOUT = 1
BLANK_SLIDE_LAYOUT = 6
TITLE_ONLY_LAYOUT = 5

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE_LAYOUT])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point

def add_image_slide(title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        # Center the image
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        txBox.text_frame.text = f"[Image {image_path} not available]"

# Slide 1: Title
add_title_slide(
    "Travel & Tour Booking System",
    "Final Exam Project\nCourse: Best Programming Practices and Design Patterns"
)

# Slide 2: Phase 1 - Problem Statement
add_bullet_slide(
    "1. Topic & Problem Statement",
    [
        "Topic: Design and Development of an Automated Travel and Tour Booking System.",
        "Case Study: Wanderlust Travels Ltd",
        "Problem Statement:",
        " - The company relies on manual booking processes using phone calls and Excel.",
        " - This leads to inefficiencies, double-bookings, and delayed payment tracking.",
        " - A centralized system is required to manage online bookings and track revenue."
    ]
)

# Render diagrams
diagrams = [
    ("Component Diagram (Functional)", "component_diagram.mermaid"),
    ("Use Case Diagram", "use_case_diagram.mermaid"),
    ("Class Diagram", "class_diagram.mermaid"),
    ("Activity Diagram", "activity diagram.mermaid"),
    ("Sequence Diagram", "sequence diagram.mermaid")
]

for title, filename in diagrams:
    filepath = os.path.join("Diagrams", filename)
    img_path = filename.replace(".mermaid", ".png")
    if os.path.exists(filepath):
        with open(filepath, "r", encoding='utf-8') as f:
            code = f.read()
        print(f"Fetching {filename}...")
        get_mermaid_image(code, img_path)
    add_image_slide(f"Phase 1: {title}", img_path)

# Slide: Phase 2 Prototype
add_bullet_slide(
    "Phase 2: Software Prototype",
    [
        "Fully functional prototype using PHP, MySQL, HTML, and CSS.",
        "User Journeys: Login -> Dashboard -> Browse Tours -> Book Tour -> Logout.",
        "Input Processing: Forms for Registration, Login, and Booking are validated.",
        "Database: Actual MySQL database integration simulates real-time availability."
    ]
)

# Slide: Phase 2 Best Practices
add_bullet_slide(
    "Phase 2: Best Programming Practices & Design Pattern",
    [
        "Practices Used:",
        " - Meaningful variable names ($tour_id vs $x).",
        " - Single Responsibility Principle separated into modular files.",
        " - Clear comments explaining logic.",
        "Design Pattern: Repository Pattern",
        " - A TourRepository class encapsulates all database logic.",
        " - Frontend calls methods like getAllTours() instead of raw SQL."
    ]
)

# Slide: Phase 3 Docker & VCS
add_bullet_slide(
    "Phase 3: Dockerizing & Version Control",
    [
        "Docker Implementation:",
        " - The app is containerized using a Dockerfile and docker-compose.yml.",
        " - Containers include Apache Web Server and MySQL Database.",
        " - Eliminates environment inconsistencies across developer machines.",
        "Version Control:",
        " - Managed using Git.",
        " - .gitignore prevents pushing sensitive credentials like the .env file."
    ]
)

# Slide: Phase 4 Test Plan
add_bullet_slide(
    "Phase 4: Software Test Plan",
    [
        "Testing Goals:",
        " - Ensure the booking system prevents double bookings and secures user data.",
        "Test Cases (Normal & Edge Cases):",
        " - TC01: User logs in with correct credentials -> Dashboard",
        " - TC02: User tries to book more slots than available -> Rejected",
        "Tools for Tracking:",
        " - PHPUnit for automated unit tests (e.g. TourRepositoryTest.php).",
        " - GitHub/GitLab Issues used as a Kanban board to track bugs."
    ]
)

prs.save("Final_Exam_Presentation.pptx")
print("Presentation generated successfully: Final_Exam_Presentation.pptx")
